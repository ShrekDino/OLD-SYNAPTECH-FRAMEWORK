#!/usr/bin/env python3
"""
SynapTechBio — Ultimate Pitch Deck Builder
Generates a 13-slide PPTX following specs from the Workfolder team documents.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# --- BRAND CONSTANTS ---
DARK_TEAL = RGBColor(0x0A, 0x1F, 0x2E)
SECTION_BG = RGBColor(0x0B, 0x2B, 0x3E)
TEAL = RGBColor(0x0D, 0x73, 0x77)
MEDIUM_TEAL = RGBColor(0x14, 0xA3, 0xA8)
GOLD = RGBColor(0xFF, 0xD1, 0x66)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xB0, 0xBE, 0xC5)
DARK_BLUE = RGBColor(0x0B, 0x2B, 0x3E)
WARM_ORANGE = RGBColor(0xF4, 0x84, 0x5F)
NEURON_GREEN = RGBColor(0x00, 0xFF, 0x88)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN_L = Inches(1.5)
MARGIN_R = Inches(1.5)
MARGIN_T = Inches(1.2)
CONTENT_W = Inches(10.333)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H


def set_bg(slide, color, gradient=False):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, size=16,
                color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                font_name="Calibri", italic=False, tracking=None):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                      Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.italic = italic
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_logo(slide):
    add_textbox(slide, 0.4, 0.35, 2.0, 0.5, "SYNAPTECHBIO",
                size=14, color=WHITE, bold=True)


def add_footer(slide, num, confidential="Confidential"):
    # Divider line
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.4), Inches(6.95), Inches(12.533), Pt(0.5))
    line.fill.solid()
    line.fill.fore_color.rgb = TEAL
    line.line.fill.background()
    line.fill.fore_color.brightness = 0.7

    slide_num = f"Slide {num:02d}"
    add_textbox(slide, 0.4, 7.0, 3.0, 0.35, slide_num,
                size=9, color=LIGHT_GRAY)
    add_textbox(slide, 3.5, 7.0, 6.0, 0.35,
                "The OS for Neuromorphic Intelligence",
                size=9, color=LIGHT_GRAY)
    add_textbox(slide, 10.5, 7.0, 2.5, 0.35, confidential,
                size=8, color=LIGHT_GRAY, italic=True)


def add_bullets(slide, left, top, width, height, items, size=14):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                      Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = LIGHT_GRAY
        p.font.name = "Calibri"
        p.space_after = Pt(6)
        p.level = 0
    return txBox


def add_bullet_block(slide, left, top, bullets, size=14, color=LIGHT_GRAY):
    return add_bullets(slide, left, top, 9, len(bullets)*0.35, bullets, size)


def add_metric_box(slide, left, top, number, label, num_color=GOLD,
                   label_color=TEAL, box_color=None):
    w, h = 2.8, 1.2
    if box_color:
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left), Inches(top), Inches(w), Inches(h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = box_color
        shape.line.fill.background()
    add_textbox(slide, left + 0.15, top + 0.1, w - 0.3, 0.55,
                str(number), size=32, color=num_color, bold=True,
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide, left + 0.15, top + 0.6, w - 0.3, 0.4,
                label, size=11, color=label_color,
                alignment=PP_ALIGN.CENTER)


def add_section_label(slide, text):
    add_textbox(slide, 1.5, 0.6, 8, 0.4, text.upper(),
                size=11, color=GOLD, bold=True, tracking=4)


# ═══════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, DARK_TEAL)

# Gradient overlay using a semi-transparent shape
overlay = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
overlay.fill.gradient()
overlay.fill.gradient_stops[0].position = 0.0
overlay.fill.gradient_stops[0].color.rgb = DARK_TEAL
overlay.fill.gradient_stops[0].color.brightness = 0.0
# Just use solid for now

add_textbox(slide, 1.5, 2.0, 10.333, 1.5, "SYNAPTECHBIO",
            size=48, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_textbox(slide, 1.5, 3.5, 10.333, 0.8,
            "The Decentralized Intelligence Foundry",
            size=24, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
add_textbox(slide, 1.5, 4.3, 10.333, 0.6,
            "From Connectome to Collective Superintelligence",
            size=16, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
add_textbox(slide, 1.5, 5.5, 10.333, 0.5,
            "Delaware C-Corp (Pre-Incorporated)",
            size=13, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
add_textbox(slide, 1.5, 6.3, 10.333, 0.4,
            "Sami Torres  |  SamiT2825@synaptechbio.org",
            size=11, color=MEDIUM_TEAL, alignment=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════
# SLIDE 2 — PROBLEM
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, DARK_TEAL)
add_logo(slide)
add_textbox(slide, 1.5, 1.5, 10, 0.7, "AI is Centralizing Power",
            size=34, color=WHITE, bold=True)
add_textbox(slide, 1.5, 2.2, 10, 0.5,
            "Faster Than We Can Regulate It",
            size=20, color=TEAL)
add_bullet_block(slide, 1.5, 3.0, [
    "$10M+ training costs for frontier models",
    "Hardware gatekept by ~200 labs worldwide (INRC)",
    "Proprietary models locked behind closed APIs",
    "Talent concentrated in 3 cities (SF, Seattle, NYC)",
    "600J per token — a fly brain runs a billion lifetimes on the same energy"
], size=14)
add_footer(slide, 2)

# ═══════════════════════════════════════════════════════════════
# SLIDE 3 — VISION: From Connectome to Computation
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, SECTION_BG)
add_logo(slide)
add_section_label(slide, "The Vision")
add_textbox(slide, 1.5, 1.5, 10, 1.0, "From Connectome to Computation",
            size=38, color=WHITE, bold=True)
add_textbox(slide, 1.5, 2.8, 10.333, 0.8,
            "The most efficient intelligence system in the universe runs on 20W. Its wiring diagram was published in Nature 2024. I built the engine that runs it.",
            size=16, color=LIGHT_GRAY)
add_metric_box(slide, 1.5, 4.0, "139,255", "Neurons", GOLD, TEAL, DARK_BLUE)
add_metric_box(slide, 4.7, 4.0, "~50M", "Synapses", GOLD, TEAL, DARK_BLUE)
add_metric_box(slide, 7.9, 4.0, "CC BY 4.0", "Open to Everyone", GOLD, GOLD, DARK_BLUE)
add_footer(slide, 3)

# ═══════════════════════════════════════════════════════════════
# SLIDE 4 — SCIENTIFIC FOUNDATION
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, DARK_TEAL)
add_logo(slide)
add_textbox(slide, 1.5, 1.5, 10, 0.7, "The Connectome Is a Gift",
            size=34, color=WHITE, bold=True)
add_textbox(slide, 1.5, 2.2, 10, 0.5,
            "FlyWire — Dorkenwald et al., Nature 2024",
            size=18, color=TEAL)

# Metric boxes
add_metric_box(slide, 1.5, 3.0, "139,255", "Neurons", GOLD, TEAL, DARK_BLUE)
add_metric_box(slide, 4.7, 3.0, "~50M", "Synapses", GOLD, TEAL, DARK_BLUE)
add_metric_box(slide, 7.9, 3.0, "CC BY 4.0", "Open License", GOLD, GOLD, DARK_BLUE)

add_bullet_block(slide, 1.5, 4.6, [
    "First complete synaptic-resolution wiring diagram of an adult brain",
    "50+ neuroscientists validated over 1,000+ hours of proofreading",
    "Eon Systems: 125k neurons driving simulated fly body in real-time — structural wiring drives behavior",
    "FlyWire LSM: >95% accuracy, trains in <30s on CPU, 1.6 MB footprint",
    "flywire-realtime-engine: 60Hz closed-loop on consumer GPU (RTX 3060)"
], size=13)
add_footer(slide, 4)

# ═══════════════════════════════════════════════════════════════
# SLIDE 5 — PRODUCT: IDRE
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, DARK_TEAL)
add_logo(slide)
add_textbox(slide, 1.5, 1.5, 10, 0.7, "The IDRE Engine",
            size=34, color=WHITE, bold=True)
add_textbox(slide, 1.5, 2.2, 10, 0.5,
            "Integrated Data Representation Engine — Real-Time Connectome Computation",
            size=18, color=TEAL)

add_bullet_block(slide, 1.5, 3.0, [
    "GPU-accelerated Compressed Sparse Column (CSC) graph engine",
    "Executes 130k × 130k Drosophila connectome — ~1ms activation on GPU, ~10ms CPU",
    "Hardware-agnostic: CuPy → SciPy → NumPy fallback — no single point of failure",
    "SSE streaming: real-time neural pulse data to browser-based 3D visualization (Three.js/R3F)",
    "Intel Loihi bridge via Lava framework — compile subgraphs for neuromorphic hardware",
    "Data Capture Split Layer (DCSL): cryptographic IP protection on every request"
], size=13)

# Performance metric row
add_metric_box(slide, 1.5, 5.5, "5x", "Latency Reduction", GOLD, GOLD, DARK_BLUE)
add_metric_box(slide, 4.7, 5.5, "20x", "Energy Savings", GOLD, GOLD, DARK_BLUE)
add_metric_box(slide, 7.9, 5.5, "25x", "Memory Efficiency", GOLD, GOLD, DARK_BLUE)
add_textbox(slide, 1.5, 6.7, 5, 0.3,
            "Stack: FastAPI  |  CuPy/SciPy  |  React 19  |  Three.js  |  Lava-NC",
            size=9, color=LIGHT_GRAY)
add_footer(slide, 5)

# ═══════════════════════════════════════════════════════════════
# SLIDE 6 — MOAT: DCSL (Section Divider)
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, SECTION_BG)
add_logo(slide)
add_section_label(slide, "The Moat")
add_textbox(slide, 1.5, 1.5, 10, 1.0, "Data Capture Split Layer",
            size=36, color=WHITE, bold=True)
add_bullet_block(slide, 1.5, 3.0, [
    "Every request forks: proprietary IP encrypted (AES-256-GCM), telemetry anonymized (SHA-256)",
    "We cannot decrypt user data — trust is engineered, not promised",
    "Creates a non-scrappable proprietary dataset that grows with every researcher",
    "Traditional AI cannot scrape what it cannot decrypt — permanent data moat",
    "Planned: per-tenant S3 buckets, Pinecone vector DB, continuous alignment pipeline"
], size=14)
add_textbox(slide, 1.5, 5.5, 10, 0.5,
            "\"They can't copy what they can't scrape.\"",
            size=18, color=GOLD, italic=True)
add_footer(slide, 6)

# ═══════════════════════════════════════════════════════════════
# SLIDE 7 — ECOSYSTEM
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, DARK_TEAL)
add_logo(slide)
add_textbox(slide, 1.5, 1.5, 10, 0.7, "The Full-Stack Ecosystem",
            size=34, color=WHITE, bold=True)

# Ecosystem cards (simplified as a structured list)
ecosystem_items = [
    ("synaptech-idre", "Core Product — Connectome computation engine", TEAL),
    ("Flywirellm", "Scientific Validation — >95% LSM accuracy", TEAL),
    ("flywire-realtime-engine", "Proof of Concept — 60Hz closed-loop simulation", MEDIUM_TEAL),
    ("EVE", "AI Alignment — Self-verifying knowledge entity", MEDIUM_TEAL),
    ("OpenMonoAgent.ai", "AI Infrastructure — We build AI tooling itself", TEAL),
    ("digidollar", "Decentralized Compute — ASIC-resistant PoW", GOLD),
    ("samchat", "Research Communication — Encrypted Matrix protocol", GOLD),
    ("the-unified-blueprint", "Organizational Design — Flat hierarchy, collective ownership", LIGHT_GRAY),
    ("Project-FreeGen", "Systems Engineering — Vulkan/FSR capability", LIGHT_GRAY),
    ("ollama-bench", "AI Benchmarking — Performance measurement infrastructure", LIGHT_GRAY)
]

y_start = 2.5
for i, (name, desc, color) in enumerate(ecosystem_items):
    row = i // 2
    col = i % 2
    x = 1.5 + col * 5.5
    y = y_start + row * 0.55
    add_textbox(slide, x, y, 2.5, 0.3, name, size=11, color=color, bold=True)
    add_textbox(slide, x + 2.6, y, 3.0, 0.3, desc, size=10, color=LIGHT_GRAY)

add_footer(slide, 7)

# ═══════════════════════════════════════════════════════════════
# SLIDE 8 — GO-TO-MARKET: THE SANDBOX
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, DARK_TEAL)
add_logo(slide)
add_textbox(slide, 1.5, 1.5, 10, 0.7, "Go-to-Market: The Sandbox",
            size=32, color=WHITE, bold=True)
add_textbox(slide, 1.5, 2.2, 10, 0.5,
            "Free Connectome Simulation. Any Browser.",
            size=18, color=TEAL)

# Sandbox mockup placeholder — centered text representing the browser UI
add_textbox(slide, 1.5, 3.0, 10, 2.5,
            "┌─────────────────────────────────────────────┐\n"
            "│       ActivationSandbox.tsx                 │\n"
            "│                                             │\n"
            "│  [ Enter text to encode → ]  [ Activate ]   │\n"
            "│                                             │\n"
            "│     ╭─────────────────────────────────╮     │\n"
            "│     │    Neuron Cascade Visualization │     │\n"
            "│     │    ▓▓▓▓▓▓░░░░ 47 Hz             │     │\n"
            "│     │    1,392 neurons active         │     │\n"
            "│     ╰─────────────────────────────────╯     │\n"
            "│                                             │\n"
            "└─────────────────────────────────────────────┘",
            size=11, color=LIGHT_GRAY)

# Sandbox URL
add_textbox(slide, 1.5, 5.7, 10, 0.4,
            "sandbox.synaptechbio.org",
            size=18, color=MEDIUM_TEAL)

# Value propositions side by side
add_textbox(slide, 1.5, 6.2, 5, 0.6,
            "Free. No auth. 100 activations/day.\nSame engine. 1.2ms spMV. Built on IDRE.",
            size=11, color=LIGHT_GRAY)
add_textbox(slide, 7.0, 6.2, 4, 0.6,
            "Enterprise: custom pricing for labs\nthat outgrow the sandbox.",
            size=11, color=LIGHT_GRAY, alignment=PP_ALIGN.RIGHT)
add_footer(slide, 8)

# ═══════════════════════════════════════════════════════════════
# SLIDE 9 — TALENT: WHY AUSTIN
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, DARK_TEAL)
add_logo(slide)
add_textbox(slide, 1.5, 1.5, 10, 0.7, "Why Austin?",
            size=34, color=WHITE, bold=True)
add_textbox(slide, 1.5, 2.2, 10, 0.5,
            "The Next Great Tech City",
            size=18, color=TEAL)

add_bullet_block(slide, 1.5, 3.0, [
    "40% lower cost of living than San Francisco",
    "No state income tax — net take-home pay is significantly higher",
    "UT Austin provides a world-class engineering pipeline",
    "Tesla, Apple, Google, Oracle — all have major Austin operations",
    "Austin doesn't have its defining tech company yet. We're building it."
], size=14)

add_textbox(slide, 1.5, 5.0, 10, 0.8,
            "\"Valve is in Bellevue. We're bringing the Valve model to Austin.\"",
            size=18, color=GOLD, italic=True)

# Talent cost comparison
add_textbox(slide, 1.5, 5.8, 5, 0.3, "Effective Cost of Engineering Talent",
            size=12, color=GOLD, bold=True)
add_textbox(slide, 1.5, 6.1, 5, 0.3,
            "SF: $495k  |  Seattle: $383k  |  NYC: $430k  |  Austin: $251k",
            size=11, color=LIGHT_GRAY)
add_textbox(slide, 7.0, 5.8, 4, 0.3, "Austin Advantage:",
            size=12, color=GOLD, bold=True)
add_textbox(slide, 7.0, 6.1, 4, 0.3,
            "49% lower effective cost than SF",
            size=11, color=NEURON_GREEN)
add_footer(slide, 9)

# ═══════════════════════════════════════════════════════════════
# SLIDE 10 — ORGANIZATION: SOLO FOUNDER
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, DARK_TEAL)
add_logo(slide)
add_textbox(slide, 1.5, 1.5, 10, 0.7, "The Organization",
            size=34, color=WHITE, bold=True)

# Centered founder text
add_textbox(slide, 1.5, 2.5, 10, 1.0, "SAMI TORRES",
            size=48, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_textbox(slide, 1.5, 3.4, 10, 0.5,
            "Solo Founder. Full Stack. Full Time.",
            size=20, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Core principles
add_bullet_block(slide, 3.0, 4.5, [
    "Decision-making: founder makes all calls",
    "Governance: ad-hoc external advisors, no formal board",
    "Ethical review: external, as needed"
], size=14)

# Future note
add_textbox(slide, 1.5, 6.0, 10, 0.4,
            "Governance model scales with team. Valve model deferred until 6+ FTE.",
            size=11, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
add_footer(slide, 10)

# ═══════════════════════════════════════════════════════════════
# SLIDE 11 — ROADMAP: 26-WEEK EXECUTION
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, DARK_TEAL)
add_logo(slide)
add_textbox(slide, 1.5, 1.5, 10, 0.7, "The Roadmap",
            size=34, color=WHITE, bold=True)
add_textbox(slide, 1.5, 2.2, 10, 0.5,
            "Solo-achievable milestones. No team dependencies.",
            size=16, color=TEAL)

phases = [
    ("P0", "Public Sandbox Deploy", "$5k", "Week 1-2", "Live sandbox URL"),
    ("P1", "STTR Partner Outreach", "$5k", "Week 3-6", "3 LOIs in hand"),
    ("P2", "Patent Filings (Batch 1)", "$20k", "Week 4-8", "4 utility patents filed"),
    ("P3", "Sandbox Iteration", "$5k", "Week 6-10", "100 DAU"),
    ("P4", "STTR Submissions (Batch 1)", "$10k", "Week 10-14", "3 grants submitted"),
    ("P5", "Loihi Benchmark Complete", "$15k", "Week 12-18", "Published results"),
    ("P6", "Fundraise Trigger", "$10k", "Week 18-26", "Next round open")
]
y = 3.0
for phase, name, cost, timeline, milestone in phases:
    add_textbox(slide, 1.5, y, 1.0, 0.3, phase, size=11, color=TEAL, bold=True)
    add_textbox(slide, 2.3, y, 2.5, 0.3, name, size=11, color=WHITE)
    add_textbox(slide, 5.0, y, 1.2, 0.3, cost, size=10, color=GOLD)
    add_textbox(slide, 6.3, y, 1.2, 0.3, timeline, size=10, color=LIGHT_GRAY)
    add_textbox(slide, 7.7, y, 4.0, 0.3, f"◆ {milestone}", size=10, color=GOLD)
    y += 0.45

add_textbox(slide, 1.5, 6.3, 10, 0.3,
            "Total: $150k  |  Every milestone is independently achievable by a solo operator.",
            size=12, color=GOLD, bold=True)
add_footer(slide, 11)

# ═══════════════════════════════════════════════════════════════
# SLIDE 12 — USE OF FUNDS: SOLO ALLOCATION
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, DARK_TEAL)
add_logo(slide)
add_textbox(slide, 1.5, 1.5, 10, 0.7, "Use of Funds: $150,000",
            size=34, color=WHITE, bold=True)
add_textbox(slide, 1.5, 2.2, 10, 0.5,
            "Zero dollars for salaries. Every dollar buys compute, patents, and proof.",
            size=14, color=GOLD)

funds = [
    ("Cloud GPU (Reserved)", "$50,000", "33%", TEAL),
    ("Legal / IP", "$35,000", "23%", MEDIUM_TEAL),
    ("INRC / Loihi Access", "$35,000", "23%", GOLD),
    ("Infrastructure", "$15,000", "10%", WARM_ORANGE),
    ("Marketing / Travel", "$15,000", "10%", LIGHT_GRAY),
    ("Engineering Salaries", "$0", "0%", DARK_BLUE)
]

y = 3.0
for label, amount, pct, color in funds:
    add_textbox(slide, 1.5, y, 4.0, 0.35, label, size=13, color=WHITE)
    add_textbox(slide, 6.0, y, 1.5, 0.35, amount, size=13, color=GOLD, bold=True)
    add_textbox(slide, 7.8, y, 1.0, 0.35, pct, size=13, color=LIGHT_GRAY)
    y += 0.5

add_textbox(slide, 1.5, 6.2, 10, 0.5,
            "$0 salaries  |  $0 office  |  $0 management  |  100% infrastructure + IP",
            size=14, color=GOLD, bold=True)
add_footer(slide, 12)

# ═══════════════════════════════════════════════════════════════
# SLIDE 13 — CLOSE / CTA
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, DARK_TEAL)

add_textbox(slide, 2.0, 2.0, 9.333, 1.5,
            "\"The future of intelligence\nis not proprietary.\nIt's collective.\"",
            size=34, color=WHITE, bold=False, italic=True,
            alignment=PP_ALIGN.CENTER)

add_textbox(slide, 2.0, 3.8, 9.333, 0.5,
            "— SynapTechBio",
            size=16, color=GOLD, alignment=PP_ALIGN.CENTER)

add_textbox(slide, 2.0, 4.8, 9.333, 0.4,
            "github.com/ShrekDino/SynapTechBio",
            size=14, color=MEDIUM_TEAL, alignment=PP_ALIGN.CENTER)

add_textbox(slide, 2.0, 5.3, 9.333, 0.4,
            "Sami Torres  |  SamiT2825@synaptechbio.org",
            size=12, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

add_textbox(slide, 2.0, 5.8, 9.333, 0.4,
            "★ Star the Repo  |  💰 Invest  |  ✉️ Connect",
            size=13, color=TEAL, alignment=PP_ALIGN.CENTER)

add_textbox(slide, 2.0, 6.5, 9.333, 0.5,
            "Delaware C-Corp (Pre-Incorporated)",
            size=10, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════════════════════════
output_path = "/home/cinni/Workfolder/SynapTechBio_Ultimate_Pitch_Deck.pptx"
prs.save(output_path)
print(f"✅ Deck saved: {output_path}")
print(f"   {len(prs.slides)} slides generated")
print(f"   Size: {os.path.getsize(output_path) / 1024:.1f} KB")
